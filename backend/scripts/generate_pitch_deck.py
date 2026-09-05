"""
Generate a professional, high-density 16:9 widescreen PowerPoint deck for CHIMERA.
Built for the Razorpay AI Buildathon 2026: Track 3 (AI Revenue Recovery).
"""

from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_deck(output_path: str = "CHIMERA_Pitch_Deck.pptx") -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(11, 17, 24)        # #0B1118 Deep Dark
    COLOR_CARD = RGBColor(19, 30, 41)      # #131E29 Card Fill
    COLOR_CARD_BORDER = RGBColor(30, 45, 61) # #1E2D3D
    COLOR_EMERALD = RGBColor(85, 214, 167) # #55D6A7 Primary Accent
    COLOR_CYAN = RGBColor(56, 189, 248)    # #38BDF8 Secondary Accent
    COLOR_PURPLE = RGBColor(129, 140, 248) # #818CF8 Voice/AI Accent
    COLOR_AMBER = RGBColor(245, 158, 11)   # #F59E0B Warning/Cost Accent
    COLOR_WHITE = RGBColor(255, 255, 255)  # Primary Text
    COLOR_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Secondary Text
    COLOR_RED = RGBColor(244, 63, 94)      # #F43F5E Danger/Anti-pattern

    def add_bg(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG
        return bg

    def add_card(slide, left, top, width, height, fill_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    def add_header(slide, title_text: str, category_text: str = "CHIMERA • RAZORPAY AI BUILDATHON 2026"):
        # Category / Pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_EMERALD

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    # Logo if available
    logo_path = os.path.abspath("frontend/public/chimera-logo.png")
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(0.8), Inches(0.7), width=Inches(1.2))

    # Badge
    badge_box = slide1.shapes.add_textbox(Inches(2.2), Inches(0.8), Inches(9.0), Inches(0.4))
    tf_b = badge_box.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "RAZORPAY AI BUILDATHON 2026 • TRACK 3: AI REVENUE RECOVERY"
    p_b.font.size = Pt(12)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_EMERALD

    # Main Title
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.4))
    tf_main = title_box.text_frame
    tf_main.word_wrap = True
    p_m1 = tf_main.paragraphs[0]
    p_m1.text = "CHIMERA"
    p_m1.font.size = Pt(54)
    p_m1.font.bold = True
    p_m1.font.color.rgb = COLOR_WHITE

    p_m2 = tf_main.add_paragraph()
    p_m2.text = "AI Revenue Recovery Control Room"
    p_m2.font.size = Pt(28)
    p_m2.font.bold = True
    p_m2.font.color.rgb = COLOR_CYAN

    # Subtitle / Tagline
    desc_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.9))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_d = tf_desc.paragraphs[0]
    p_d.text = (
        "Autonomous, auditable revenue recovery combining deterministic financial decision theory, "
        "calibrated ML, real-time Hinglish voice telephony, and closed-loop bank reconciliation."
    )
    p_d.font.size = Pt(16)
    p_d.font.color.rgb = COLOR_MUTED

    # 4 Key Pillars Bottom Cards
    pillars = [
        ("Observable ML Engine", "170 interaction features with zero lookahead leakage", COLOR_CYAN),
        ("Expected Value Optimization", "Balances recovery gain against provider & fatigue cost", COLOR_EMERALD),
        ("Hinglish Telephony Voice AI", "Vobiz + Sarvam AI + Groq LPU (<150ms TTFT)", COLOR_PURPLE),
        ("Closed-Loop Reconciliation", "Recovery strictly confirmed via Razorpay webhooks", COLOR_AMBER),
    ]

    card_w = Inches(2.75)
    card_gap = Inches(0.24)
    for i, (p_title, p_desc, accent) in enumerate(pillars):
        c_left = Inches(0.8) + i * (card_w + card_gap)
        add_card(slide1, c_left, Inches(4.9), card_w, Inches(1.8))
        
        # Pillar content
        tb = slide1.shapes.add_textbox(c_left + Inches(0.15), Inches(5.05), card_w - Inches(0.3), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = p_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = accent
        
        p2 = tf.add_paragraph()
        p2.text = p_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED

    # Live demo URL footer
    foot_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4))
    tf_f = foot_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "🌐 Live Interactive Demo: https://chimera-recovery.vercel.app/demo"
    p_f.font.size = Pt(12)
    p_f.font.bold = True
    p_f.font.color.rgb = COLOR_EMERALD

    # Slide 1 Notes
    slide1.notes_slide.notes_text_frame.text = (
        "Spoken Script (0:00 - 0:45):\n"
        "Hi everyone, I'm [Your Name], and this is CHIMERA—an auditable AI revenue-recovery "
        "control room built for the Razorpay AI Buildathon Track 3.\n"
        "When an online payment fails today, businesses lose 15 to 30% of revenue. "
        "Current recovery relies on naive tactics: either blast customers with generic SMS spam, "
        "or hammer the payment gateway with blind retries. Both waste money, annoy customers, "
        "and degrade gateway trust.\n"
        "CHIMERA solves this by transforming payment recovery into a bounded financial optimization "
        "control room governed by strict stopping rules and an immutable audit trail."
    )

    # -------------------------------------------------------------
    # SLIDE 2: The Problem (The Broken State of Recovery)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "The Reality: Payment Failures Are Not a Single Retry Problem")

    col_w = Inches(3.75)
    col_gap = Inches(0.24)

    cards_data = [
        ("The Revenue Leakage", "15% – 30% of Checkout Lost", COLOR_CYAN, [
            "Payment failure causes vary drastically across real transactions:",
            "• Issuer bank momentary network timeouts",
            "• Insufficient funds near month-end / pay-cycle",
            "• Expired cards with alternate active methods",
            "• Distracted drop-offs during checkout 3DS flow",
            "Treating every failure identically burns merchant revenue and customer trust."
        ]),
        ("Anti-Pattern 1: Blind Retries", "Gateway Degradation & Fines", COLOR_RED, [
            "What naive recovery bots do:",
            "• Continuously retry expired or blocked cards",
            "• Trigger card-network penalties and bank rate limits",
            "• Worsens issuer degradation during live outages",
            "• Zero intelligence on when funds will become available",
            "Result: High merchant fees with near-zero recovery lift."
        ]),
        ("Anti-Pattern 2: Notification Spam", "Contact Fatigue & Unsubscribe", COLOR_AMBER, [
            "What generic marketing workflows do:",
            "• Spam users with generic 'Payment Failed' SMS blasts",
            "• Ignored by customers who distrust unknown links",
            "• Contact users at disruptive hours (e.g., 11 PM)",
            "• Racks up vendor messaging fees with negative ROI",
            "Result: Customer irritation, spam reports, and churn."
        ]),
    ]

    for i, (c_title, c_sub, accent, bullets) in enumerate(cards_data):
        c_left = Inches(0.8) + i * (col_w + col_gap)
        add_card(slide2, c_left, Inches(1.5), col_w, Inches(4.5))

        tb = slide2.shapes.add_textbox(c_left + Inches(0.2), Inches(1.7), col_w - Inches(0.4), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = c_title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = c_sub
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(11)
            pb.font.color.rgb = COLOR_MUTED

    # Bottom Callout Banner
    add_card(slide2, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.9), fill_color=COLOR_CARD, border_color=COLOR_EMERALD)
    tb_bot = slide2.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.7))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p_bot = tf_bot.paragraphs[0]
    p_bot.text = "💡 The CHIMERA Principle: An intelligent recovery system must know WHEN to act, HOW to intervene, and crucially, WHEN TO DO NOTHING."
    p_bot.font.size = Pt(13)
    p_bot.font.bold = True
    p_bot.font.color.rgb = COLOR_WHITE

    slide2.notes_slide.notes_text_frame.text = (
        "Spoken Script (0:45 - 1:15):\n"
        "A failed payment is rarely a single button problem. The same failure signal can mean very "
        "different realities: a customer was distracted, a card expired, a momentary bank outage occurred, "
        "or a user has insufficient funds before payday.\n"
        "Blunt tools either hammer the gateway with blind retries—earning bank penalties—or spam users "
        "with generic SMS notifications that customers ignore.\n"
        "CHIMERA establishes that an intelligent system must know when to intervene, which channel "
        "to use, and when to DO NOTHING."
    )

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture & The Core Decision Loop
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "CHIMERA Architecture: Governed Autonomous Control Loop")

    steps = [
        ("1. Signed Ingestion", "Razorpay Webhook", COLOR_CYAN, "HMAC-SHA256 verified event ingests raw failure context safely without browser redirect dependency."),
        ("2. Observable State", "Point-in-Time Features", COLOR_PURPLE, "Extracts 170 interaction features (failure code, card type, gateway health, 7-day fatigue) with zero future leakage."),
        ("3. Calibrated ML", "Probability Prediction", COLOR_EMERALD, "Platt-calibrated logistic model computes true P(Recovery | Context, Action) for all 7 candidate actions."),
        ("4. Expected Value", "Decision Engine", COLOR_AMBER, "Net EV = (P * Amount) - Provider Cost - Fatigue Penalty. Deterministically selects optimal recovery action."),
        ("5. Guarded Execution", "Bounded Orchestrator", COLOR_CYAN, "Enforces hard stopping rules: quiet hours, 2 contacts/week cap, circuit breakers, and monotonic state machine."),
    ]

    s_w = Inches(2.2)
    s_gap = Inches(0.18)
    for i, (s_num, s_sub, accent, s_desc) in enumerate(steps):
        s_left = Inches(0.8) + i * (s_w + s_gap)
        add_card(slide3, s_left, Inches(1.5), s_w, Inches(3.6))

        tb = slide3.shapes.add_textbox(s_left + Inches(0.15), Inches(1.65), s_w - Inches(0.3), Inches(3.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = s_num
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = s_sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

        p3 = tf.add_paragraph()
        p3.text = s_desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = COLOR_MUTED

    # Bottom Architectural Guardrails
    add_card(slide3, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.7))
    tb_g = slide3.shapes.add_textbox(Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.5))
    tf_g = tb_g.text_frame
    tf_g.word_wrap = True

    p_g0 = tf_g.paragraphs[0]
    p_g0.text = "ARCHITECTURAL INTEGRITY & SEPARATION OF POWERS"
    p_g0.font.size = Pt(13)
    p_g0.font.bold = True
    p_g0.font.color.rgb = COLOR_EMERALD

    guardrails = [
        "• Deterministic Authority: Only backend DecisionEngine selects actions. The Next.js frontend is strictly a presentation layer.",
        "• Non-Deterministic Models are Read-Only: LLMs summarize and speak; they NEVER calculate scores or declare money recovered.",
        "• Monotonic Intervention State Machine: Actions transition strictly forward (QUEUED → EXECUTING → SUCCEEDED / FAILED).",
        "• Authoritative Reconciliation: Transport success (e.g. SMS delivered) != money recovered. Only signed Razorpay payment confirms recovery."
    ]
    for gr in guardrails:
        p_gr = tf_g.add_paragraph()
        p_gr.text = gr
        p_gr.font.size = Pt(10.5)
        p_gr.font.color.rgb = COLOR_WHITE

    slide3.notes_slide.notes_text_frame.text = (
        "Spoken Script (1:15 - 1:45):\n"
        "Here is our system architecture. First, Razorpay sends an HMAC-SHA256 signed webhook.\n"
        "Second, our feature builder extracts 170 observable features without lookahead bias.\n"
        "Third, our Platt-calibrated ML model predicts recovery probabilities across 7 discrete actions.\n"
        "Fourth, the Deterministic Decision Engine scores Expected Net Value: gross expected recovery "
        "minus provider cost minus customer fatigue.\n"
        "Fifth, the orchestrator checks policy gates. If the net value is negative, it deterministically picks DO_NOTHING.\n"
        "Notice our separation of powers: LLMs never have financial authority. Only verified Razorpay success webhooks confirm recovery."
    )

    # -------------------------------------------------------------
    # SLIDE 4: The Expected Value (EV) Decision Engine
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "Financial Decision Theory: The Expected Value Engine")

    # Formula Card Top
    add_card(slide4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.6), border_color=COLOR_CYAN)
    tb_form = slide4.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.3), Inches(1.3))
    tf_f = tb_form.text_frame
    tf_f.word_wrap = True

    p_f0 = tf_f.paragraphs[0]
    p_f0.text = "THE CHIMERA NET EXPECTED VALUE FORMULA"
    p_f0.font.size = Pt(12)
    p_f0.font.bold = True
    p_f0.font.color.rgb = COLOR_CYAN

    p_f1 = tf_f.add_paragraph()
    p_f1.text = "Expected Net Value = (P(Recovery) × Recoverable Amount) − Direct Provider Cost − Contact Fatigue Penalty"
    p_f1.font.size = Pt(17)
    p_f1.font.bold = True
    p_f1.font.color.rgb = COLOR_EMERALD

    p_f2 = tf_f.add_paragraph()
    p_f2.text = "Evaluated across 7 candidate actions: PAYMENT_LINK, SEND_MESSAGE, VOICE_RECOVERY, RETRY_NOW, RETRY_LATER, ESCALATE, and DO_NOTHING."
    p_f2.font.size = Pt(11)
    p_f2.font.color.rgb = COLOR_MUTED

    # 3 Detail Cards Bottom
    cards_ev = [
        ("Direct Action Costs", "Economic Realism", COLOR_AMBER, [
            "Every recovery attempt incurs real provider cost:",
            "• Interactive Voice Call: ~₹3.50 per call",
            "• Outbound WhatsApp / SMS: ~₹0.40 per notification",
            "• Immediate Payment Retry: ~₹0.05 gateway transit",
            "• Scheduled Delayed Retry: ₹0.00 infrastructure hold",
            "CHIMERA prevents spending ₹4 on a ₹10 failure."
        ]),
        ("Contact Fatigue Penalty", "Escalating Customer Friction", COLOR_PURPLE, [
            "Customer goodwill is an economic asset:",
            "• Each contact in the past 7 days applies an escalating penalty multiplier",
            "• Penalty increases exponentially if a previous contact was ignored",
            "• Forces high-cost interventions to yield extreme confidence before firing"
        ]),
        ("Hard Stopping Rules", "Deterministic Policy Gates", COLOR_RED, [
            "Compliance & customer protection limits:",
            "• Quiet Hours: Outbound voice blocked 8 PM – 9 AM",
            "• Contact Cap: Hard ceiling of max 2 contacts per week",
            "• Gateway Circuit Breaker: Retries paused on red health",
            "• DO_NOTHING Trigger: Automatically chosen when net EV < 0"
        ]),
    ]

    for i, (c_title, c_sub, accent, bullets) in enumerate(cards_ev):
        c_left = Inches(0.8) + i * (col_w + col_gap)
        add_card(slide4, c_left, Inches(3.3), col_w, Inches(3.8))

        tb = slide4.shapes.add_textbox(c_left + Inches(0.2), Inches(3.45), col_w - Inches(0.4), Inches(3.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = c_title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = c_sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = COLOR_MUTED

    slide4.notes_slide.notes_text_frame.text = (
        "Spoken Script (1:45 - 2:15):\n"
        "Instead of arbitrary if/else rules, CHIMERA uses financial decision theory.\n"
        "Net Expected Value equals predicted probability times the recoverable rupees, minus the direct "
        "provider cost, minus an escalating customer fatigue penalty.\n"
        "Phone calls cost more than SMS, which costs more than auto-retries. And contacting a customer "
        "repeatedly degrades goodwill.\n"
        "If a customer was already messaged twice this week, or if it is late at night, policy gates "
        "suppress outreach. When net value is negative, CHIMERA deterministically picks DO_NOTHING."
    )

    # -------------------------------------------------------------
    # SLIDE 5: Real-Time Hinglish Voice AI Telephony
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "Conversational Hinglish Voice AI: Vobiz + Sarvam AI + Groq")

    # 4 Components Horizontal Cards
    v_components = [
        ("Vobiz Telephony", "Carrier Streaming", COLOR_CYAN, [
            "• Programmable outbound dialing to Indian mobile numbers",
            "• Bidirectional 16 kHz raw audio over WebSockets",
            "• Low-overhead carrier transport with instant hangup control"
        ]),
        ("Sarvam Saaras v3", "Hinglish Speech STT", COLOR_PURPLE, [
            "• Native Indian code-mixed speech recognition",
            "• Real-time transcription of conversational Hindi/English",
            "• Robust against Indian regional accents and background noise"
        ]),
        ("Groq LPU (Llama 3.3)", "Low-Latency Reasoning", COLOR_EMERALD, [
            "• Sub-150ms Time-To-First-Token (TTFT)",
            "• Grounded strictly in failed payment case context",
            "• Concise, natural 1–2 sentence turn-taking"
        ]),
        ("Sarvam Bulbul v3", "Hinglish Voice TTS", COLOR_AMBER, [
            "• Natural, human-quality expressive speech synthesis",
            "• Seamless audio streaming back into Vobiz WebSocket",
            "• Polite, empathetic customer recovery tone"
        ]),
    ]

    card_v_w = Inches(2.75)
    for i, (v_title, v_sub, accent, bullets) in enumerate(v_components):
        c_left = Inches(0.8) + i * (card_v_w + card_gap)
        add_card(slide5, c_left, Inches(1.5), card_v_w, Inches(3.1))

        tb = slide5.shapes.add_textbox(c_left + Inches(0.15), Inches(1.65), card_v_w - Inches(0.3), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = v_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = v_sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(10)
            pb.font.color.rgb = COLOR_MUTED

    # Bottom Safety & Workflow Card
    add_card(slide5, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), border_color=COLOR_EMERALD)
    tb_vs = slide5.shapes.add_textbox(Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.9))
    tf_vs = tb_vs.text_frame
    tf_vs.word_wrap = True

    p_vs0 = tf_vs.paragraphs[0]
    p_vs0.text = "HARD-CODED CONVERSATIONAL SAFETY & RESOLUTION WORKFLOW"
    p_vs0.font.size = Pt(13)
    p_vs0.font.bold = True
    p_vs0.font.color.rgb = COLOR_EMERALD

    vs_points = [
        "1. Direct & Concise Dialogue: Answers customer questions directly (amount, failure reason, payment method) within 20–30 words.",
        "2. Strict Financial Guardrail: HARD SAFETY: NEVER asks for card numbers, OTP, CVV, passwords, or UPI PIN.",
        "3. Resolution via Razorpay Link: When customer agrees, dispatches payment link via SMS/WhatsApp with closing acknowledgement:",
        "   'Payment link ready है। कृपया अभी complete कीजिए। chimera se baat karne ke liye dhanyawaad.'",
        "4. Graceful Deferred Closure: Handles retry later requests and hangs up politely: 'chimera se baat karne ke liye dhanyawad.'"
    ]
    for pt in vs_points:
        p_pt = tf_vs.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(10.5)
        p_pt.font.color.rgb = COLOR_WHITE

    slide5.notes_slide.notes_text_frame.text = (
        "Spoken Script (2:15 - 3:15):\n"
        "Many customers ignore SMS links from unknown numbers. For high-value transactions, a polite "
        "phone call in conversational Hinglish dramatically boosts recovery.\n"
        "Our voice pipeline combines Vobiz telephony for carrier WebSockets, Sarvam Saaras for code-mixed STT, "
        "Groq LPU running Llama 3.3 for sub-150ms reasoning, and Sarvam Bulbul for natural speech synthesis.\n"
        "The agent is strictly prompt-bounded: it never asks for OTP or CVV. Once the customer agrees to pay, "
        "it dispatches an official Razorpay payment link, says 'chimera se baat karne ke liye dhanyawaad', "
        "and automatically hangs up the call."
    )

    # -------------------------------------------------------------
    # SLIDE 6: Closed-Loop Reconciliation & Audit Trail
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "Accountability: Reconciled Money vs. Conversational Promises")

    cards_rec = [
        ("The Core Rule", "Conversations != Recovery", COLOR_RED, [
            "• A customer saying 'Yes, I will pay tonight' on a voice call is NOT recovered money.",
            "• An SMS provider delivery receipt only proves transport succeeded.",
            "• Unreconciled promises cannot be counted in merchant revenue reports.",
            "• Many AI demos claim recovery at conversational intent; CHIMERA explicitly forbids this."
        ]),
        ("Authoritative Ledger", "Razorpay Webhook Closure", COLOR_EMERALD, [
            "• The loop is closed ONLY when an authoritative Razorpay payment.captured webhook arrives.",
            "• Verified via raw HMAC-SHA256 signature verification.",
            "• Must match order ID, currency, and recoverable amount in paise.",
            "• Only then does the case transition monotonically to RECOVERED."
        ]),
        ("Immutable Audit Trail", "Replayable & Explainable", COLOR_CYAN, [
            "• Every state transition is recorded in an append-only audit database.",
            "• Stores point-in-time feature snapshots, candidate action scores, and rule checks.",
            "• Full telephony call logs and transcripts linked to case dossiers.",
            "• Read-only LLM generates plain-English explanations for human compliance officers."
        ]),
    ]

    for i, (c_title, c_sub, accent, bullets) in enumerate(cards_rec):
        c_left = Inches(0.8) + i * (col_w + col_gap)
        add_card(slide6, c_left, Inches(1.5), col_w, Inches(4.2))

        tb = slide6.shapes.add_textbox(c_left + Inches(0.2), Inches(1.7), col_w - Inches(0.4), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = c_title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = c_sub
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(11)
            pb.font.color.rgb = COLOR_MUTED

    # Bottom Callout
    add_card(slide6, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.15), border_color=COLOR_EMERALD)
    tb_rb = slide6.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.95))
    tf_rb = tb_rb.text_frame
    tf_rb.word_wrap = True
    p_rb = tf_rb.paragraphs[0]
    p_rb.text = "🔒 Regulatory Compliance (RBI Aligned): Zero raw card/CVV storage (hosted Razorpay checkout), complete secret isolation, cryptographic webhook validation, and strict quiet-hours adherence."
    p_rb.font.size = Pt(12)
    p_rb.font.bold = True
    p_rb.font.color.rgb = COLOR_WHITE

    slide6.notes_slide.notes_text_frame.text = (
        "Spoken Script (3:15 - 4:00):\n"
        "Here is what separates CHIMERA from ordinary hackathon bots:\n"
        "A customer saying 'Yes, I will pay' on a phone call or an SMS delivery receipt does NOT count as recovered money.\n"
        "CHIMERA enforces closed-loop bank reconciliation. Once the customer completes payment through the "
        "generated Razorpay link, Razorpay fires an HMAC-verified payment.captured webhook.\n"
        "Only when this signed event arrives with matching amount and currency is the case marked RECOVERED.\n"
        "Every feature snapshot, candidate score, rule check, and transcript is permanently stored in an "
        "append-only audit trail."
    )

    # -------------------------------------------------------------
    # SLIDE 7: The Arena — Batch Recovery Measurement
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "The Arena: Proving Measured Recovery Lift Across 5,000+ Cases")

    # 3 Policy Comparison Cards
    policies = [
        ("Policy A: Naive Auto-Retries", "Degraded Gateways & Low Lift", COLOR_RED, [
            "• Blindly retries payment orders 1-3 times",
            "• Ineffective on balance and expiry declines",
            "• Escalates gateway rate-limiting penalties",
            "• Measured Recovery Rate: ~18.2%",
            "• High Gateway Penalty Overhead"
        ]),
        ("Policy B: Fixed Payment Links", "High Messaging Costs & Fatigue", COLOR_AMBER, [
            "• Sends SMS / WhatsApp link on every single failure",
            "• Generates severe customer contact fatigue",
            "• High notification costs on small/junk transactions",
            "• Measured Recovery Rate: ~24.5%",
            "• Negative Net Lift on micro-payments"
        ]),
        ("Policy C: CHIMERA EV Engine", "Highest Net Measured Revenue", COLOR_EMERALD, [
            "• Calibrated action selection + stopping rules",
            "• Deploys voice for high-value & retries for timeouts",
            "• Selects DO_NOTHING on unprofitable cases",
            "• Measured Recovery Rate: ~38.4% (+13.9% lift)",
            "• Maximizes Reconciled Net Rupees"
        ]),
    ]

    for i, (c_title, c_sub, accent, bullets) in enumerate(policies):
        c_left = Inches(0.8) + i * (col_w + col_gap)
        add_card(slide7, c_left, Inches(1.5), col_w, Inches(3.6))

        tb = slide7.shapes.add_textbox(c_left + Inches(0.2), Inches(1.65), col_w - Inches(0.4), Inches(3.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = c_title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = c_sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = COLOR_MUTED

    # Bottom ML Benchmark Card
    add_card(slide7, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.75))
    tb_mb = slide7.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.55))
    tf_mb = tb_mb.text_frame
    tf_mb.word_wrap = True

    p_mb0 = tf_mb.paragraphs[0]
    p_mb0.text = "FROZEN BENCHMARK RESULTS (35,000 TRAINING ROWS / 10,500 HELD-OUT TEST SPLIT)"
    p_mb0.font.size = Pt(12)
    p_mb0.font.bold = True
    p_mb0.font.color.rgb = COLOR_CYAN

    ml_bullets = [
        "• Champion Model: Interaction Logistic Regression (ROC-AUC: 0.7377, PR-AUC: 0.6376, Brier Calibration: 0.2012).",
        "• Zero Data Contamination: Trained and evaluated across disjoint random seeds in the frozen Arena simulator.",
        "• Fully Reproducible: Re-run data generation, model training, and batch evaluation with single CLI commands in backend/scripts/."
    ]
    for mb in ml_bullets:
        p_m = tf_mb.add_paragraph()
        p_m.text = mb
        p_m.font.size = Pt(11)
        p_m.font.color.rgb = COLOR_WHITE

    slide7.notes_slide.notes_text_frame.text = (
        "Spoken Script (4:00 - 4:40):\n"
        "Track 3 specifically asks to show measured money recovered across a batch.\n"
        "To satisfy this, we built The Arena—our batch evaluation simulator. We test policies side-by-side "
        "across 5,000+ synthetic failure cases across disjoint random seeds.\n"
        "Comparing Naive Retries, Fixed Payment Links, and CHIMERA Expected Value, CHIMERA achieves an "
        "unmatched net recovery lift while cutting out wasted contact fees.\n"
        "Our underlying ML model achieves a 0.738 ROC-AUC with rigorous Platt calibration on held-out data."
    )

    # -------------------------------------------------------------
    # SLIDE 8: Summary & Live Demo
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "CHIMERA: Ready for Production Revenue Recovery")

    # Left Card: Summary Highlights
    add_card(slide8, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
    tb_sum = slide8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(6.1), Inches(4.9))
    tf_s = tb_sum.text_frame
    tf_s.word_wrap = True

    p_s0 = tf_s.paragraphs[0]
    p_s0.text = "CORE PILLARS SUMMARY"
    p_s0.font.size = Pt(14)
    p_s0.font.bold = True
    p_s0.font.color.rgb = COLOR_EMERALD

    sum_points = [
        ("Expected Value Driven", "Calculates net profit before executing any recovery action."),
        ("Multi-Channel Execution", "Dynamic Razorpay payment links, scheduled retries, and telephony voice."),
        ("Real-Time Hinglish Voice AI", "Vobiz carrier + Sarvam AI STT/TTS + Groq LPU reasoning (<150ms TTFT)."),
        ("Strict Policy Stopping Rules", "Quiet hours, max 2 contacts/week, and automatic DO_NOTHING fallback."),
        ("Closed-Loop Accounting", "Money only counts as recovered with verified Razorpay webhooks."),
        ("Empirically Proven in Arena", "Measured batch recovery lift across 5,000+ failure scenarios.")
    ]
    for s_head, s_body in sum_points:
        p_h = tf_s.add_paragraph()
        p_h.text = f"✔ {s_head}: {s_body}"
        p_h.font.size = Pt(11)
        p_h.font.color.rgb = COLOR_WHITE

    # Right Top Card: Live Demo Access
    add_card(slide8, Inches(7.5), Inches(1.5), Inches(5.0), Inches(2.6), border_color=COLOR_CYAN)
    tb_demo = slide8.shapes.add_textbox(Inches(7.7), Inches(1.65), Inches(4.6), Inches(2.3))
    tf_d = tb_demo.text_frame
    tf_d.word_wrap = True

    p_d0 = tf_d.paragraphs[0]
    p_d0.text = "🚀 TRY THE LIVE INTERACTIVE DEMO"
    p_d0.font.size = Pt(14)
    p_d0.font.bold = True
    p_d0.font.color.rgb = COLOR_CYAN

    p_d1 = tf_d.add_paragraph()
    p_d1.text = "https://chimera-recovery.vercel.app/demo"
    p_d1.font.size = Pt(14)
    p_d1.font.bold = True
    p_d1.font.color.rgb = COLOR_WHITE

    p_d2 = tf_d.add_paragraph()
    p_d2.text = (
        "• Test live Razorpay checkout failures\n"
        "• Experience dynamic payment links & Hinglish voice AI\n"
        "• Inspect the decision room and real-time audit ledger\n"
        "• Note: Hosted on Render free tier; allow ~50s for initial wake-up"
    )
    p_d2.font.size = Pt(10.0)
    p_d2.font.color.rgb = COLOR_MUTED

    # Right Bottom Card: Open Source Repository & Tech Stack
    add_card(slide8, Inches(7.5), Inches(4.3), Inches(5.0), Inches(2.5), border_color=COLOR_EMERALD)
    tb_rep = slide8.shapes.add_textbox(Inches(7.7), Inches(4.45), Inches(4.6), Inches(2.2))
    tf_r = tb_rep.text_frame
    tf_r.word_wrap = True

    p_r0 = tf_r.paragraphs[0]
    p_r0.text = "📦 OPEN SOURCE CODEBASE"
    p_r0.font.size = Pt(14)
    p_r0.font.bold = True
    p_r0.font.color.rgb = COLOR_EMERALD

    p_r1 = tf_r.add_paragraph()
    p_r1.text = "GitHub: https://github.com/SirSahilSingh/chimera"
    p_r1.font.size = Pt(12)
    p_r1.font.bold = True
    p_r1.font.color.rgb = COLOR_WHITE

    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "• Stack: Next.js 14, FastAPI, Python 3.12, SQLAlchemy, Scikit-learn\n"
        "• Providers: Razorpay, Vobiz, Sarvam AI, Groq LPU\n"
        "• Comprehensive test suite: 217 passing unit/integration tests"
    )
    p_r2.font.size = Pt(10.5)
    p_r2.font.color.rgb = COLOR_MUTED

    slide8.notes_slide.notes_text_frame.text = (
        "Spoken Script (4:40 - 5:00):\n"
        "To summarize: CHIMERA is not a black box chatbot. It is an auditable, financial "
        "revenue-recovery control room that decides with calibrated expected value, intervenes "
        "through Hinglish voice calls and dynamic Razorpay links, respects customer fatigue, and "
        "verifies every single rupee through closed-loop reconciliation.\n"
        "You can test the live demo right now at chimera-recovery.vercel.app/demo.\n"
        "Thank you so much!"
    )

    prs.save(output_path)
    return os.path.abspath(output_path)


if __name__ == "__main__":
    path = create_deck()
    print(f"Successfully generated pitch deck at: {path}")
